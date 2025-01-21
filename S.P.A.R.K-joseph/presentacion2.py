# This Python file uses the following encoding: utf-8 
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Country Don: A Presentation"
subtitle.text = "Exploring the Fictional Nation"


#slide 2: Geography
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Geography of Country Don"
tf = body_shape.text_frame
tf.text = "Country Don is a fictional nation, so its geography can be as imaginative as you like.  Perhaps it's a mountainous region, a sprawling desert, or a collection of islands. The possibilities are endless!"


#slide 3: Culture
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

shapes.title.text = "Culture of Country Don"
tf = shapes.placeholders[1].text_frame
tf.text = "Describe the culture of Country Don. What are its traditions, customs, and values? What kind of art, music, and literature does it produce?"


#slide 4: History
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

shapes.title.text = "History of Country Don"
tf = shapes.placeholders[1].text_frame
tf.text = "Craft a history for Country Don.  What significant events have shaped its development? Who are its important figures?"


#slide 5: Economy
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

shapes.title.text = "Economy of Country Don"
tf = shapes.placeholders[1].text_frame
tf.text = "Detail the economic system of Country Don. What are its major industries? What is its level of development?"


#slide 6: Politics
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

shapes.title.text = "Politics of Country Don"
tf = shapes.placeholders[1].text_frame
tf.text = "Explain the political system of Country Don. Is it a democracy, a monarchy, or something else? Who holds power?"


#Save the presentation
prs.save("country_don.pptx")

#Open the presentation (platform dependent)
os.startfile("country_don.pptx") #For Windows
#subprocess.call(('open', "country_don.pptx")) #For macOS

